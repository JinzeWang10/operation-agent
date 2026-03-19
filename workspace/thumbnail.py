"""Generate thumbnail images from a PPTX file for visual verification."""

import subprocess
import sys
import os


def pptx_to_images_libreoffice(pptx_path, output_dir):
    """Convert PPTX to images using LibreOffice."""
    os.makedirs(output_dir, exist_ok=True)

    # First convert to PDF
    cmd = [
        'soffice', '--headless', '--convert-to', 'pdf',
        '--outdir', output_dir, pptx_path
    ]
    print(f'Running: {" ".join(cmd)}')
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            print(f'LibreOffice error: {result.stderr}')
            return False
        print(f'PDF created in {output_dir}')
        return True
    except FileNotFoundError:
        print('LibreOffice (soffice) not found on PATH.')
        return False


def pptx_to_images_python(pptx_path, output_dir):
    """Convert PPTX to images using python-pptx + Pillow (basic text rendering)."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print('Install Pillow: pip install Pillow')
        return False

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Scale to reasonable image size
    img_w = 1280
    img_h = int(img_w * Emu(slide_h) / Emu(slide_w))
    scale = img_w / Emu(slide_w)

    for idx, slide in enumerate(prs.slides):
        img = Image.new('RGB', (img_w, img_h), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            x = int(shape.left * scale) if shape.left else 0
            y = int(shape.top * scale) if shape.top else 0
            w = int(shape.width * scale) if shape.width else 0
            h = int(shape.height * scale) if shape.height else 0

            # Draw shape fill
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    color = shape.fill.fore_color.rgb
                    r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                    draw.rectangle([x, y, x + w, y + h], fill=(r, g, b))
                except Exception:
                    draw.rectangle([x, y, x + w, y + h], outline=(200, 200, 200))

            # Draw text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        try:
                            font = ImageFont.truetype("msyh.ttc", max(8, int(12 * scale * 50)))
                        except Exception:
                            font = ImageFont.load_default()
                        text_color = (0, 0, 0)
                        if para.runs:
                            try:
                                c = para.runs[0].font.color.rgb
                                text_color = (int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                            except Exception:
                                pass
                        draw.text((x + 4, y + 2), text, fill=text_color, font=font)
                        y += 16

        out_path = os.path.join(output_dir, f'slide_{idx + 1:02d}.png')
        img.save(out_path)
        print(f'Saved: {out_path}')

    return True


if __name__ == '__main__':
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_dir = os.path.dirname(script_dir)
    pptx_path = os.path.join(project_dir, 'docs', 'operation-agent-proposal.pptx')
    output_dir = os.path.join(script_dir, 'slides')

    if not os.path.exists(pptx_path):
        print(f'File not found: {pptx_path}')
        sys.exit(1)

    print(f'Converting: {pptx_path}')
    print(f'Output dir: {output_dir}')

    # Try LibreOffice first
    if pptx_to_images_libreoffice(pptx_path, output_dir):
        print('\nDone! Check the output directory for PDF/images.')
    else:
        print('\nLibreOffice not available. Trying Python-based rendering...')
        if pptx_to_images_python(pptx_path, output_dir):
            print('\nDone! Basic thumbnail images generated.')
        else:
            print('\nFailed to generate thumbnails.')
            sys.exit(1)
