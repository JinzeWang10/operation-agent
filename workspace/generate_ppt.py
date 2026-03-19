"""Generate operation-agent-proposal.pptx from PICC template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ── Colors ──────────────────────────────────────────────────────────────
BLUE = RGBColor(0x5B, 0x9B, 0xD5)
DARK_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_BLUE_BG = RGBColor(0xDE, 0xEB, 0xF7)
LIGHT_ORANGE_BG = RGBColor(0xFD, 0xE9, 0xD9)
GREEN = RGBColor(0x70, 0xAD, 0x47)
RED = RGBColor(0xC0, 0x00, 0x00)

FONT_CN = '微软雅黑'
FONT_EN = 'Arial'

# ── Helpers ─────────────────────────────────────────────────────────────

def set_font(run, size=14, bold=False, color=BLACK, name_cn=FONT_CN, name_en=FONT_EN):
    """Set font properties on a run."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name_en
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name_cn)


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, word_wrap=True):
    """Add a simple text box with one paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=font_size, bold=bold, color=color)
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))
    return txBox


def add_rich_textbox(slide, left, top, width, height, paragraphs_data,
                     anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple styled paragraphs.
    paragraphs_data: list of dicts with keys:
      runs: list of (text, size, bold, color)
      alignment: PP_ALIGN value
      space_before: Pt value (optional)
      space_after: Pt value (optional)
      bullet: bool (optional)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))

    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = pdata.get('alignment', PP_ALIGN.LEFT)
        if 'space_before' in pdata:
            p.space_before = pdata['space_before']
        if 'space_after' in pdata:
            p.space_after = pdata['space_after']
        if 'line_spacing' in pdata:
            p.line_spacing = pdata['line_spacing']
        if pdata.get('bullet'):
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
            pPr.append(buChar)
            p.level = pdata.get('level', 0)
            if p.level == 0:
                pPr.set('marL', str(Pt(18).emu))
                pPr.set('indent', str(-Pt(18).emu))
            else:
                pPr.set('marL', str(Pt(36).emu))
                pPr.set('indent', str(-Pt(18).emu))

        for run_data in pdata.get('runs', []):
            text, size, bold, color = run_data
            run = p.add_run()
            run.text = text
            set_font(run, size=size, bold=bold, color=color)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BLUE,
                     text='', font_size=12, font_color=WHITE, bold=True):
    """Add a rounded rectangle shape with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Reduce corner radius
    shape.adjustments[0] = 0.1
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        shape.text_frame._txBody.bodyPr.set('anchor', 'ctr')
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        set_font(run, size=font_size, bold=bold, color=font_color)
    return shape


def add_rect(slide, left, top, width, height, fill_color=BLUE, line_color=None):
    """Add a plain rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_arrow_right(slide, left, top, width, height, fill_color=GRAY):
    """Add a right arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title_text, page_num=None):
    """Add a colored title bar at the top of the slide."""
    # Title background bar
    bar = add_rect(slide, Inches(0), Inches(0.55), Inches(13.333), Inches(0.65),
                   fill_color=DARK_BLUE)
    # Title text
    add_textbox(slide, Inches(0.5), Inches(0.58), Inches(10), Inches(0.6),
                title_text, font_size=24, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    if page_num:
        add_textbox(slide, Inches(12), Inches(0.58), Inches(0.8), Inches(0.6),
                    str(page_num), font_size=12, color=WHITE,
                    alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_tag(slide, left, top, text, fill_color=BLUE, font_size=10, width=None):
    """Add a small tag/badge."""
    w = width or Inches(1.2)
    shape = add_rounded_rect(slide, left, top, w, Inches(0.3),
                             fill_color=fill_color, text=text,
                             font_size=font_size, font_color=WHITE, bold=True)
    return shape


def add_table_to_slide(slide, left, top, width, height, rows, cols):
    """Add a table and return it."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def set_cell(table, row, col, text, font_size=11, bold=False, color=BLACK,
             fill_color=None, alignment=PP_ALIGN.LEFT):
    """Set text in a table cell."""
    cell = table.cell(row, col)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=font_size, bold=bold, color=color)


# ── Slide Builders ──────────────────────────────────────────────────────

def build_slide1_cover(slide):
    """Slide 1: Cover page."""
    # Semi-transparent dark overlay for title area
    overlay = add_rect(slide, Inches(0.8), Inches(2.2), Inches(8.5), Inches(3),
                       fill_color=DARK_BLUE)
    # Make it semi-transparent via XML
    spPr = overlay._element.spPr
    solidFill_elem = spPr.find(qn('a:solidFill'))
    if solidFill_elem is not None:
        srgbClr = solidFill_elem.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '85000'})
            srgbClr.append(alpha)

    # Main title
    add_textbox(slide, Inches(1.2), Inches(2.5), Inches(7.5), Inches(1.2),
                '运维巡检 AI Agent 方案', font_size=40, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    # Subtitle
    add_textbox(slide, Inches(1.2), Inches(3.7), Inches(7.5), Inches(0.6),
                '用 AI 替代人工巡检，提升事件响应效率', font_size=20,
                bold=False, color=RGBColor(0xDD, 0xDD, 0xDD),
                anchor=MSO_ANCHOR.TOP)
    # Divider line
    line = add_rect(slide, Inches(1.2), Inches(3.55), Inches(2.5), Inches(0.04),
                    fill_color=ORANGE)
    # Date / department
    add_textbox(slide, Inches(1.2), Inches(4.4), Inches(7.5), Inches(0.4),
                '信息技术部  |  2025', font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))


def build_slide2_pain_points(slide):
    """Slide 2: Current status and pain points."""
    add_title_bar(slide, '现状与痛点', 2)

    # Left: pain point description
    pain_items = [
        ('事件响应流程低效', [
            '事件发生后，运维人员需手动逐一登录 4-5 套监控系统',
            '单次巡检耗时 15-30 分钟',
            '依赖个人经验，检查范围因人而异，容易遗漏关键信息',
        ]),
        ('多系统割裂', [
            'Prometheus、Zabbix、ELK、云监控各自独立',
            '无法一键获取全局视图',
            '不同系统的告警缺乏关联分析',
        ]),
        ('效率瓶颈', [
            '夜间 / 节假日响应不及时',
            '重复劳动消耗大量人力',
            '经验难以沉淀和传承',
        ]),
    ]

    y_pos = Inches(1.5)
    for title, items in pain_items:
        # Section title with orange accent
        add_rect(slide, Inches(0.5), y_pos, Inches(0.08), Inches(0.35),
                 fill_color=ORANGE)
        add_textbox(slide, Inches(0.75), y_pos - Pt(2), Inches(5), Inches(0.4),
                    title, font_size=16, bold=True, color=DARK_BLUE)
        for j, item in enumerate(items):
            add_textbox(slide, Inches(1.0), y_pos + Inches(0.4 + j * 0.32),
                        Inches(5.5), Inches(0.32),
                        '\u2022  ' + item, font_size=12, color=BLACK)
        y_pos += Inches(0.4 + len(items) * 0.32 + 0.2)

    # Right: illustrative flow showing the chaos
    rx = Inches(7.2)
    add_textbox(slide, rx, Inches(1.5), Inches(5.5), Inches(0.4),
                '当前巡检流程', font_size=14, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)

    # Flow: Event -> Manual steps
    box_w = Inches(2.2)
    box_h = Inches(0.5)
    add_rounded_rect(slide, rx + Inches(1.5), Inches(2.1), box_w, box_h,
                     fill_color=RED, text='事件告警', font_size=12)
    add_arrow_right(slide, rx + Inches(1.5) + box_w + Inches(0.1),
                    Inches(2.15), Inches(0.5), Inches(0.4), fill_color=GRAY)

    # Multiple monitoring systems
    systems = ['Prometheus', 'Zabbix', 'ELK', '云监控', '其他系统']
    sy = Inches(2.9)
    for i, sys_name in enumerate(systems):
        add_rounded_rect(slide, rx + Inches(0.4), sy + Inches(i * 0.55),
                         Inches(2.0), Inches(0.42),
                         fill_color=RGBColor(0x99, 0x99, 0x99),
                         text=sys_name, font_size=11)
    # Brace-like text
    add_textbox(slide, rx + Inches(2.6), sy + Inches(0.5), Inches(3), Inches(2),
                '人工逐一登录\n逐一检查\n逐一记录\n\n= 15~30 分钟',
                font_size=13, bold=True, color=RED)


def build_slide3_solution_overview(slide):
    """Slide 3: Solution overview."""
    add_title_bar(slide, '解决方案概览', 3)

    # One-liner
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.9),
             fill_color=LIGHT_BLUE_BG)
    add_textbox(slide, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.8),
                '自然语言描述事件  →  Agent 自动并行扫描所有监控系统  →  生成巡检报告推送到事件群',
                font_size=18, bold=True, color=DARK_BLUE,
                anchor=MSO_ANCHOR.MIDDLE)

    # Three key points in rounded boxes
    points = [
        ('类比', '7×24 不知疲倦的\n巡检专家', BLUE),
        ('核心能力', '自动扫描 + AI 分析\n生成结构化报告', ORANGE),
        ('明确边界', '只做异常发现\n不做根因定位', RGBColor(0x70, 0xAD, 0x47)),
    ]
    box_w = Inches(3.5)
    box_h = Inches(1.8)
    start_x = Inches(0.8)
    gap = Inches(0.45)
    y = Inches(3.0)

    for i, (label, desc, color) in enumerate(points):
        x = start_x + i * (box_w + gap)
        # Card background
        card = add_rounded_rect(slide, x, y, box_w, box_h, fill_color=WHITE)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        card.shadow.inherit = False
        # Label tag
        add_tag(slide, x + Inches(0.2), y + Inches(0.2), label,
                fill_color=color, width=Inches(1.3))
        # Description
        add_textbox(slide, x + Inches(0.3), y + Inches(0.7), box_w - Inches(0.6),
                    Inches(1.0), desc, font_size=16, color=BLACK,
                    anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
                '* 后续可扩展：根因分析、自动修复建议等高级能力',
                font_size=11, color=GRAY)


def build_slide4_architecture(slide):
    """Slide 4: Architecture diagram — dual triggers, metadata config, conditional LLM."""
    add_title_bar(slide, '整体架构图', 4)

    # ── Left column: Two trigger sources ──
    # Trigger A: 事件工单系统
    add_rounded_rect(slide, Inches(0.3), Inches(1.7), Inches(2.0), Inches(0.65),
                     fill_color=ORANGE, text='事件工单系统',
                     font_size=12, bold=True)
    add_textbox(slide, Inches(0.3), Inches(2.4), Inches(2.0), Inches(0.3),
                '工单建立 → 自动 API 调用', font_size=8, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # Trigger B: 事件群 @Bot
    add_rounded_rect(slide, Inches(0.3), Inches(3.0), Inches(2.0), Inches(0.65),
                     fill_color=ORANGE, text='事件群 @Bot',
                     font_size=12, bold=True)
    add_textbox(slide, Inches(0.3), Inches(3.7), Inches(2.0), Inches(0.3),
                '@巡检bot 自然语言描述', font_size=8, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # Arrows from triggers to agent
    add_arrow_right(slide, Inches(2.4), Inches(1.88), Inches(0.45), Inches(0.35),
                    fill_color=BLUE)
    add_arrow_right(slide, Inches(2.4), Inches(3.15), Inches(0.45), Inches(0.35),
                    fill_color=BLUE)

    # ── Center: Agent Service ──
    agent_x, agent_y = Inches(3.0), Inches(1.4)
    agent_w, agent_h = Inches(4.5), Inches(4.8)
    agent_box = add_rect(slide, agent_x, agent_y, agent_w, agent_h,
                         fill_color=WHITE, line_color=DARK_BLUE)
    agent_box.line.width = Pt(2)
    add_textbox(slide, agent_x, agent_y + Inches(0.08), agent_w, Inches(0.35),
                'Agent 服务（三阶段编排引擎）', font_size=14, bold=True,
                color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # 输入规范化 step
    add_rounded_rect(slide, agent_x + Inches(0.2), agent_y + Inches(0.5),
                     Inches(4.1), Inches(0.42),
                     fill_color=RGBColor(0x88, 0x88, 0x88),
                     text='输入规范化 → IncidentRequest',
                     font_size=11)

    # Phase 1 (with conditional LLM annotation)
    p1_y = agent_y + Inches(1.1)
    add_rounded_rect(slide, agent_x + Inches(0.2), p1_y,
                     Inches(4.1), Inches(0.95), fill_color=BLUE,
                     text='Phase 1: 基础巡检\n并行拉取 + 多维粗筛',
                     font_size=11)
    add_textbox(slide, agent_x + Inches(0.3), p1_y + Inches(0.68),
                Inches(3.9), Inches(0.22),
                '* 告警 > 阈值时条件触发 LLM 精筛', font_size=8,
                color=RGBColor(0xCC, 0xDD, 0xFF))

    # Phase 2
    p2_y = agent_y + Inches(2.2)
    add_rounded_rect(slide, agent_x + Inches(0.2), p2_y,
                     Inches(4.1), Inches(0.95), fill_color=ORANGE,
                     text='Phase 2: AI 深入调查\nTool Calling 查指标/日志',
                     font_size=11)

    # Phase 3
    p3_y = agent_y + Inches(3.3)
    add_rounded_rect(slide, agent_x + Inches(0.2), p3_y,
                     Inches(4.1), Inches(0.95), fill_color=GREEN,
                     text='Phase 3: 报告生成\nAI 生成 + 降级模板兜底',
                     font_size=11)

    # ── Right top: Monitoring systems ──
    mon_x = Inches(8.1)
    add_textbox(slide, mon_x, Inches(1.4), Inches(3), Inches(0.35),
                '监控系统', font_size=13, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)
    systems = ['Prometheus', 'Zabbix', 'ELK', '云监控']
    for i, sys_name in enumerate(systems):
        add_rounded_rect(slide, mon_x + Inches(0.2), Inches(1.85 + i * 0.55),
                         Inches(2.2), Inches(0.42),
                         fill_color=BLUE, text=sys_name, font_size=11)

    # Adapter label
    add_textbox(slide, Inches(7.6), Inches(2.3), Inches(0.6), Inches(1.0),
                '适\n配\n器', font_size=10, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_arrow_right(slide, Inches(7.6), Inches(2.1), Inches(0.4), Inches(0.35),
                    fill_color=BLUE)

    # ── Right middle: 元数据配置 ──
    cfg_y = Inches(4.2)
    add_rounded_rect(slide, mon_x + Inches(0.2), cfg_y, Inches(2.2),
                     Inches(0.5), fill_color=RGBColor(0x88, 0x88, 0x88),
                     text='元数据配置', font_size=11)
    add_textbox(slide, mon_x, cfg_y + Inches(0.55), Inches(2.6), Inches(0.25),
                '指标名 / PromQL 示例 / 索引名', font_size=8, color=GRAY,
                alignment=PP_ALIGN.CENTER)
    # Arrow from config to Agent (pointing left)
    add_textbox(slide, Inches(7.6), cfg_y + Inches(0.08), Inches(0.5), Inches(0.35),
                '\u2190', font_size=16, bold=True, color=RGBColor(0x88, 0x88, 0x88),
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── Right bottom: LLM Service ──
    llm_y = Inches(5.1)
    add_textbox(slide, mon_x, llm_y, Inches(3), Inches(0.3),
                '内部 LLM 服务', font_size=13, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, mon_x + Inches(0.2), llm_y + Inches(0.35),
                     Inches(2.2), Inches(0.45), fill_color=ORANGE,
                     text='OpenAI 兼容接口', font_size=11)
    add_textbox(slide, Inches(7.6), llm_y + Inches(0.15), Inches(0.6), Inches(0.5),
                '抽\n象\n层', font_size=10, bold=True, color=ORANGE,
                alignment=PP_ALIGN.CENTER)
    add_arrow_right(slide, Inches(7.6), llm_y + Inches(0.4), Inches(0.4), Inches(0.35),
                    fill_color=ORANGE)

    # ── Bottom left: Output ──
    add_rounded_rect(slide, Inches(0.3), Inches(4.5), Inches(2.0), Inches(0.55),
                     fill_color=GREEN, text='Webhook 推送',
                     font_size=11, bold=True)
    add_textbox(slide, Inches(0.3), Inches(5.1), Inches(2.0), Inches(0.3),
                '钉钉 / 飞书 / Slack', font_size=9, color=GRAY,
                alignment=PP_ALIGN.CENTER)
    add_arrow_right(slide, Inches(2.4), Inches(4.6), Inches(0.45), Inches(0.35),
                    fill_color=GREEN)

    # ── Footer note ──
    add_rect(slide, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.45),
             fill_color=LIGHT_BLUE_BG)
    add_textbox(slide, Inches(0.5), Inches(6.32), Inches(12), Inches(0.4),
                '全部组件内网部署  |  技术栈全开源  |  数据不出内网',
                font_size=13, bold=True, color=DARK_BLUE,
                anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)


def build_slide5_three_phases(slide):
    """Slide 5: Three-phase workflow — Code vs AI responsibilities (core page)."""
    add_title_bar(slide, '三阶段工作流 — 代码 vs AI 职责', 5)

    # Tag: core page
    add_tag(slide, Inches(10.5), Inches(0.6), '核心页', fill_color=ORANGE,
            width=Inches(0.9), font_size=10)

    # Phase 1
    p1_x, p1_y = Inches(0.3), Inches(1.5)
    p1_w, p1_h = Inches(3.9), Inches(5.0)
    card1 = add_rect(slide, p1_x, p1_y, p1_w, p1_h, fill_color=WHITE,
                     line_color=BLUE)
    card1.line.width = Pt(2)
    add_rounded_rect(slide, p1_x + Inches(0.15), p1_y + Inches(0.15),
                     Inches(3.6), Inches(0.55), fill_color=BLUE,
                     text='Phase 1  基础巡检', font_size=15)
    add_tag(slide, p1_x + Inches(0.15), p1_y + Inches(0.85), '代码控制',
            fill_color=RGBColor(0x00, 0x70, 0xC0), width=Inches(1.2), font_size=10)
    add_tag(slide, p1_x + Inches(1.5), p1_y + Inches(0.85), '可选 AI',
            fill_color=RGBColor(0xFF, 0x99, 0x33), width=Inches(1.0), font_size=10)

    p1_items = [
        ('1.1 并行拉取活跃告警', False),
        ('1.2 多维粗筛', False),
        ('时间 / 系统名 / 严重级 / 关键词', True),
        ('1.3 条件 LLM 精筛（>阈值时触发）', False),
        ('告警过多 \u2192 LLM 判断相关性', True),
        ('输出：优先级分组的告警列表', False),
    ]
    for i, (item, is_sub) in enumerate(p1_items):
        left = p1_x + (Inches(0.55) if is_sub else Inches(0.3))
        prefix = '  \u2023  ' if is_sub else '\u2022  '
        fsize = 10 if is_sub else 11
        fcolor = GRAY if is_sub else BLACK
        add_textbox(slide, left, p1_y + Inches(1.3 + i * 0.35),
                    Inches(3.1), Inches(0.32),
                    prefix + item, font_size=fsize, color=fcolor)

    # Output box
    add_rect(slide, p1_x + Inches(0.2), p1_y + Inches(3.7), Inches(3.5),
             Inches(0.55), fill_color=LIGHT_BLUE_BG)
    add_textbox(slide, p1_x + Inches(0.3), p1_y + Inches(3.72), Inches(3.3),
                Inches(0.5), 'BaselineScanResult\nalerts + alerts_by_priority',
                font_size=9, color=DARK_BLUE)

    # Decorative progress bar
    add_rect(slide, p1_x + Inches(0.3), p1_y + Inches(4.5), Inches(3.3),
             Inches(0.18), fill_color=BLUE)

    # Phase 2
    p2_x = Inches(4.5)
    card2 = add_rect(slide, p2_x, p1_y, p1_w, p1_h, fill_color=WHITE,
                     line_color=ORANGE)
    card2.line.width = Pt(2)
    add_rounded_rect(slide, p2_x + Inches(0.15), p1_y + Inches(0.15),
                     Inches(3.6), Inches(0.55), fill_color=ORANGE,
                     text='Phase 2  深入调查', font_size=15)
    add_tag(slide, p2_x + Inches(0.15), p1_y + Inches(0.85), 'AI 决策',
            fill_color=ORANGE, width=Inches(1.0), font_size=10)
    add_tag(slide, p2_x + Inches(1.3), p1_y + Inches(0.85), '代码执行',
            fill_color=RGBColor(0x00, 0x70, 0xC0), width=Inches(1.2), font_size=10)

    p2_items = [
        ('上下文组装（代码）', False),
        ('事件 + 告警 + 系统元数据', True),
        ('LLM 写原生查询语句', False),
        ('PromQL / ES DSL', True),
        ('1 tool call = 1 轮，最多 3 轮', False),
        ('查询失败 \u2192 错误回传 LLM', False),
        ('超时/轮次耗尽 \u2192 代码强制结束', False),
    ]
    for i, (item, is_sub) in enumerate(p2_items):
        left = p2_x + (Inches(0.55) if is_sub else Inches(0.3))
        prefix = '  \u2023  ' if is_sub else '\u2022  '
        fsize = 10 if is_sub else 11
        fcolor = GRAY if is_sub else BLACK
        add_textbox(slide, left, p1_y + Inches(1.3 + i * 0.33),
                    Inches(3.1), Inches(0.3),
                    prefix + item, font_size=fsize, color=fcolor)

    # Output box
    add_rect(slide, p2_x + Inches(0.2), p1_y + Inches(3.7), Inches(3.5),
             Inches(0.55), fill_color=LIGHT_ORANGE_BG)
    add_textbox(slide, p2_x + Inches(0.3), p1_y + Inches(3.72), Inches(3.3),
                Inches(0.5), 'InvestigationResult\nfindings + queries_executed',
                font_size=9, color=ORANGE)

    add_rect(slide, p2_x + Inches(0.3), p1_y + Inches(4.5), Inches(3.3),
             Inches(0.18), fill_color=ORANGE)

    # Phase 3
    p3_x = Inches(8.7)
    card3 = add_rect(slide, p3_x, p1_y, p1_w, p1_h, fill_color=WHITE,
                     line_color=GREEN)
    card3.line.width = Pt(2)
    add_rounded_rect(slide, p3_x + Inches(0.15), p1_y + Inches(0.15),
                     Inches(3.6), Inches(0.55), fill_color=GREEN,
                     text='Phase 3  报告生成', font_size=15)
    add_tag(slide, p3_x + Inches(0.15), p1_y + Inches(0.85), 'AI 生成',
            fill_color=GREEN, width=Inches(1.0), font_size=10)
    add_tag(slide, p3_x + Inches(1.3), p1_y + Inches(0.85), '代码推送',
            fill_color=RGBColor(0x00, 0x70, 0xC0), width=Inches(1.2), font_size=10)

    p3_items = [
        ('AI 生成 Markdown 结构化报告', False),
        ('降级策略', False),
        ('LLM 不可用 \u2192 模板引擎生成', True),
        ('代码格式化 + Webhook 推送', False),
        ('钉钉 / 飞书 / Slack', True),
    ]
    for i, (item, is_sub) in enumerate(p3_items):
        left = p3_x + (Inches(0.55) if is_sub else Inches(0.3))
        prefix = '  \u2023  ' if is_sub else '\u2022  '
        fsize = 10 if is_sub else 11
        fcolor = GRAY if is_sub else BLACK
        add_textbox(slide, left, p1_y + Inches(1.3 + i * 0.35),
                    Inches(3.1), Inches(0.32),
                    prefix + item, font_size=fsize, color=fcolor)

    # Output box
    add_rect(slide, p3_x + Inches(0.2), p1_y + Inches(3.7), Inches(3.5),
             Inches(0.55), fill_color=RGBColor(0xE2, 0xEF, 0xDA))
    add_textbox(slide, p3_x + Inches(0.3), p1_y + Inches(3.72), Inches(3.3),
                Inches(0.5), 'InspectionReport\nreport_markdown + notification',
                font_size=9, color=GREEN)

    add_rect(slide, p3_x + Inches(0.3), p1_y + Inches(4.5), Inches(3.3),
             Inches(0.18), fill_color=GREEN)

    # Arrows between phases
    add_arrow_right(slide, Inches(4.15), Inches(3.2), Inches(0.35), Inches(0.4),
                    fill_color=GRAY)
    add_arrow_right(slide, Inches(8.35), Inches(3.2), Inches(0.35), Inches(0.4),
                    fill_color=GRAY)


def build_slide6_agent_detail(slide):
    """Slide 6: Agent work process detail — 10-step flow (key page)."""
    add_title_bar(slide, 'Agent 工作原理详解', 6)

    # Scenario
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.55),
             fill_color=LIGHT_ORANGE_BG)
    add_textbox(slide, Inches(0.7), Inches(1.42), Inches(11.5), Inches(0.5),
                '示例："事件工单 \u2014 XX 业务访问超时"',
                font_size=14, bold=True, color=ORANGE, anchor=MSO_ANCHOR.MIDDLE)

    # Steps flow - vertical timeline (10 steps)
    steps = [
        ('[代码]',     '接收请求，规范化输入（fault_system=订单服务, event_time=14:30）'),
        ('[代码]',     '并行查询 Prometheus / Zabbix / ELK，拉取 14:30\u00b160min 内告警'),
        ('[代码]',     '多维粗筛：系统名"订单服务" + severity\u2265warning \u2192 45\u219218条'),
        ('[代码]',     '18条 \u2264 阈值30，跳过 LLM 精筛，直接进入 Phase 2'),
        ('[代码]',     '组装上下文：事件描述 + 18条告警 + Prometheus 元数据（指标名列表）'),
        ('[AI\u2192Tool]', '第1轮：query_metrics("prometheus", "rate(http_requests_total{...}[5m])")'),
        ('[代码]',     '执行 PromQL，返回结果给 AI'),
        ('[AI\u2192Tool]', '第2轮：search_logs("elasticsearch", {"query":{"match":{...}}})'),
        ('[  AI  ]',   '调用 finish_investigation \u2014 信息充足，结束调查'),
        ('[AI+代码]',  'AI 生成报告 \u2192 代码格式化 Markdown \u2192 Webhook 推送到事件群'),
    ]

    # Phase boundary info: (start_step_idx, end_step_idx, phase_name, color)
    phase_ranges = [
        (0, 4, 'Phase 1', BLUE),
        (5, 8, 'Phase 2', ORANGE),
        (9, 9, 'Phase 3', GREEN),
    ]

    start_y = Inches(2.15)
    step_h = Inches(0.34)
    gap = Inches(0.05)

    for i, (tag, desc) in enumerate(steps):
        y = start_y + i * (step_h + gap)

        # Determine color from phase ranges
        color = BLUE
        for ps, pe, _, pc in phase_ranges:
            if ps <= i <= pe:
                color = pc
                break

        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y,
                                        Inches(0.32), Inches(0.32))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        circle.text_frame._txBody.bodyPr.set('anchor', 'ctr')
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        set_font(run, size=9, bold=True, color=WHITE)

        # Vertical line connector (except last)
        if i < len(steps) - 1:
            add_rect(slide, Inches(0.64), y + Inches(0.32),
                     Inches(0.04), gap, fill_color=RGBColor(0xCC, 0xCC, 0xCC))

        # Tag
        tag_color = RGBColor(0x00, 0x70, 0xC0) if color == BLUE else (GREEN if color == GREEN else ORANGE)
        tag_w = Inches(1.2) if len(tag) > 6 else Inches(0.9)
        add_rounded_rect(slide, Inches(1.0), y + Inches(0.02), tag_w,
                         Inches(0.28), fill_color=tag_color, text=tag,
                         font_size=9, bold=True)

        # Description
        desc_x = Inches(1.0) + tag_w + Inches(0.15)
        desc_w = Inches(12.5) - desc_x
        add_textbox(slide, desc_x, y, desc_w, step_h, desc,
                    font_size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)

    # Phase boundary labels on the right
    for ps, pe, phase_name, color in phase_ranges:
        y_start = start_y + ps * (step_h + gap)
        y_end = start_y + pe * (step_h + gap) + step_h
        bar_h = y_end - y_start
        mid_y = y_start + (bar_h - Inches(0.25)) // 2
        # Vertical color bar
        add_rect(slide, Inches(12.7), y_start, Inches(0.06), bar_h,
                 fill_color=color)
        # Phase label
        add_textbox(slide, Inches(12.0), mid_y, Inches(0.7), Inches(0.25),
                    phase_name, font_size=8, bold=True, color=color,
                    alignment=PP_ALIGN.RIGHT)

    # Summary box at bottom
    add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             fill_color=LIGHT_BLUE_BG)
    add_textbox(slide, Inches(0.7), Inches(6.22), Inches(11.5), Inches(0.45),
                '全程自动化  |  代码保证确定性  |  AI 提供智能决策  |  总耗时 1-2 分钟',
                font_size=13, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_slide7_interaction(slide):
    """Slide 7: How to interact with the Agent — dual trigger sources."""
    add_title_bar(slide, '如何与 Agent 交互', 7)

    # ── Input: Trigger Source A ──
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(1.5),
             fill_color=WHITE, line_color=BLUE)
    add_tag(slide, Inches(0.7), Inches(1.6), '触发源 A', fill_color=BLUE,
            width=Inches(1.1))
    items_a = [
        ('事件工单系统 \u2014 工单建立时自动调用 API', 12, BLACK),
        ('携带结构化字段：fault_system / description / alert_summary', 10, GRAY),
        ('event_time / time_window_minutes（默认 60）', 10, GRAY),
    ]
    for i, (item, sz, clr) in enumerate(items_a):
        add_textbox(slide, Inches(0.8), Inches(2.05 + i * 0.32),
                    Inches(5.2), Inches(0.3), item, font_size=sz, color=clr)

    # ── Input: Trigger Source B ──
    add_rect(slide, Inches(0.5), Inches(3.2), Inches(5.8), Inches(1.2),
             fill_color=WHITE, line_color=BLUE)
    add_tag(slide, Inches(0.7), Inches(3.3), '触发源 B', fill_color=BLUE,
            width=Inches(1.1))
    items_b = [
        ('事件群 @巡检bot \u2014 发送自然语言描述', 12, BLACK),
        ('系统解析为相同 IncidentRequest 结构（一期单次巡检）', 10, GRAY),
    ]
    for i, (item, sz, clr) in enumerate(items_b):
        add_textbox(slide, Inches(0.8), Inches(3.75 + i * 0.32),
                    Inches(5.2), Inches(0.3), item, font_size=sz, color=clr)

    # Arrow
    add_arrow_right(slide, Inches(6.5), Inches(2.7), Inches(0.6), Inches(0.5),
                    fill_color=DARK_BLUE)

    # ── Output section ──
    add_rect(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(2.9),
             fill_color=WHITE, line_color=GREEN)
    add_tag(slide, Inches(7.5), Inches(1.6), '输出', fill_color=GREEN,
            width=Inches(0.8))
    items_output = [
        'Markdown 格式巡检报告',
        '自动推送到事件群（Webhook）',
        '包含：告警列表、指标异常、日志摘要',
        '同时通过 API 返回给调用方',
    ]
    for i, item in enumerate(items_output):
        add_textbox(slide, Inches(7.5), Inches(2.1 + i * 0.4),
                    Inches(5.2), Inches(0.38), item, font_size=12, color=BLACK)

    # ── Future expansion ──
    add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8),
             fill_color=LIGHT_GRAY, line_color=GRAY)
    add_tag(slide, Inches(0.7), Inches(4.9), '后续扩展', fill_color=GRAY,
            width=Inches(1.2))
    future_items = [
        '支持多轮追问和对话式交互（一期只做单次巡检）',
        '支持批量巡检和定时巡检',
        '更多监控系统适配器接入',
    ]
    for i, item in enumerate(future_items):
        add_textbox(slide, Inches(0.8), Inches(5.35 + i * 0.4),
                    Inches(11.5), Inches(0.38),
                    '\u2022  ' + item, font_size=13, color=BLACK)


def build_slide8_adapters(slide):
    """Slide 8: Plugin-based adapter architecture."""
    add_title_bar(slide, '监控系统对接 — 插件化适配器', 8)

    # Unified interface concept
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
                '每个监控系统一个适配器，统一接口规范',
                font_size=16, bold=True, color=DARK_BLUE)

    # Adapter cards
    adapters = [
        ('Prometheus', '指标查询\n活跃告警', BLUE, '已规划'),
        ('Zabbix', '主机监控\n触发器告警', BLUE, '已规划'),
        ('ELK', '日志搜索\n错误聚合', BLUE, '已规划'),
        ('CloudWatch', 'AWS 云监控\n指标 & 告警', ORANGE, '已规划'),
        ('阿里云', '阿里云监控\n云产品指标', ORANGE, '已规划'),
    ]

    card_w = Inches(2.2)
    card_h = Inches(2.0)
    start_x = Inches(0.5)
    gap = Inches(0.3)
    y = Inches(2.2)

    for i, (name, desc, color, status) in enumerate(adapters):
        x = start_x + i * (card_w + gap)
        card = add_rect(slide, x, y, card_w, card_h, fill_color=WHITE,
                        line_color=color)
        card.line.width = Pt(1.5)
        # Name
        add_rounded_rect(slide, x + Inches(0.1), y + Inches(0.15),
                         card_w - Inches(0.2), Inches(0.45),
                         fill_color=color, text=name, font_size=13)
        # Description
        add_textbox(slide, x + Inches(0.15), y + Inches(0.75),
                    card_w - Inches(0.3), Inches(0.8),
                    desc, font_size=11, color=GRAY,
                    alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Status tag
        add_tag(slide, x + Inches(0.3), y + Inches(1.6), status,
                fill_color=color, width=Inches(1.0), font_size=9)

    # Extension method
    add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.8),
             fill_color=LIGHT_BLUE_BG)
    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
                '扩展方式', font_size=15, bold=True, color=DARK_BLUE)

    ext_steps = [
        ('1', '拿到接口文档'),
        ('2', '开发适配器（实现统一接口）'),
        ('3', '加配置项'),
        ('4', '接入完成'),
    ]
    for i, (num, text) in enumerate(ext_steps):
        sx = Inches(0.8) + i * Inches(3.0)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx, Inches(5.4),
                                        Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = DARK_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        circle.text_frame._txBody.bodyPr.set('anchor', 'ctr')
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, size=12, bold=True, color=WHITE)

        add_textbox(slide, sx + Inches(0.5), Inches(5.38), Inches(2.3),
                    Inches(0.45), text, font_size=13, color=BLACK,
                    anchor=MSO_ANCHOR.MIDDLE)
        if i < len(ext_steps) - 1:
            add_arrow_right(slide, sx + Inches(2.5), Inches(5.45),
                            Inches(0.35), Inches(0.3), fill_color=BLUE)

    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
                '核心编排逻辑不需要改动，适配器即插即用',
                font_size=12, bold=True, color=ORANGE)


def build_slide9_deployment(slide):
    """Slide 9: Intranet deployment feasibility."""
    add_title_bar(slide, '内网部署可行性', 9)

    # Tech stack
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.5),
             fill_color=WHITE, line_color=BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.55), Inches(5), Inches(0.45),
                '技术栈（全开源）', font_size=16, bold=True, color=DARK_BLUE)

    stack_items = [
        ('Python 3.11+', '主语言'),
        ('FastAPI', 'Web 框架'),
        ('httpx', '异步 HTTP 客户端'),
        ('Pydantic', '数据校验'),
        ('structlog', '结构化日志'),
        ('Docker', '容器化部署'),
    ]
    for i, (tech, desc) in enumerate(stack_items):
        y = Inches(2.2 + i * 0.42)
        add_textbox(slide, Inches(0.9), y, Inches(2.2), Inches(0.38),
                    tech, font_size=13, bold=True, color=DARK_BLUE)
        add_textbox(slide, Inches(3.2), y, Inches(2.8), Inches(0.38),
                    desc, font_size=12, color=GRAY)

    # LLM integration
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(1.8),
             fill_color=WHITE, line_color=ORANGE)
    add_textbox(slide, Inches(7.0), Inches(1.55), Inches(5), Inches(0.45),
                'LLM 对接方案', font_size=16, bold=True, color=ORANGE)
    llm_items = [
        '通过抽象层对接，兼容任何 OpenAI 兼容接口',
        '内部 LLM 服务提供接口后即可对接',
        '支持灵活切换模型',
    ]
    for i, item in enumerate(llm_items):
        add_textbox(slide, Inches(7.0), Inches(2.2 + i * 0.38),
                    Inches(5.3), Inches(0.35),
                    '\u2022  ' + item, font_size=12, color=BLACK)

    # Security / network
    add_rect(slide, Inches(6.8), Inches(3.7), Inches(5.8), Inches(1.3),
             fill_color=WHITE, line_color=GREEN)
    add_textbox(slide, Inches(7.0), Inches(3.75), Inches(5), Inches(0.4),
                '网络与安全', font_size=16, bold=True, color=GREEN)
    sec_items = [
        '无需访问外网，所有数据流转在内网完成',
        'Docker 容器化，轻量级部署',
    ]
    for i, item in enumerate(sec_items):
        add_textbox(slide, Inches(7.0), Inches(4.3 + i * 0.35),
                    Inches(5.3), Inches(0.32),
                    '\u2022  ' + item, font_size=12, color=BLACK)

    # Deployment diagram
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.3),
             fill_color=LIGHT_BLUE_BG)
    add_textbox(slide, Inches(0.7), Inches(5.45), Inches(12), Inches(0.4),
                '部署架构', font_size=14, bold=True, color=DARK_BLUE)

    deploy_boxes = [
        'Docker 容器', 'FastAPI 服务', '适配器模块', 'LLM 客户端', 'Webhook 推送'
    ]
    for i, label in enumerate(deploy_boxes):
        bx = Inches(0.7 + i * 2.5)
        add_rounded_rect(slide, bx, Inches(5.95), Inches(2.0), Inches(0.5),
                         fill_color=DARK_BLUE, text=label, font_size=11)
        if i < len(deploy_boxes) - 1:
            add_arrow_right(slide, bx + Inches(2.05), Inches(6.05),
                            Inches(0.35), Inches(0.3), fill_color=BLUE)


def build_slide10_comparison(slide):
    """Slide 10: Before/After comparison + sample report."""
    add_title_bar(slide, '效果对比 + 使用场景', 10)

    # Comparison table
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(5), Inches(0.4),
                '对比分析', font_size=16, bold=True, color=DARK_BLUE)

    table = add_table_to_slide(slide, Inches(0.5), Inches(1.9),
                               Inches(6.5), Inches(3.6), 5, 3)
    # Set column widths
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)

    # Header row
    headers = [('维度', WHITE), ('人工巡检', WHITE), ('AI Agent', WHITE)]
    for j, (text, color) in enumerate(headers):
        set_cell(table, 0, j, text, font_size=13, bold=True, color=color,
                 fill_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # Data rows
    rows_data = [
        ('耗时', '15-30 分钟', '1-2 分钟'),
        ('覆盖范围', '依赖个人经验', '全系统无遗漏'),
        ('一致性', '因人而异', '标准化输出'),
        ('可用性', '工作时间', '7×24 小时'),
    ]
    for i, (dim, before, after) in enumerate(rows_data):
        row_idx = i + 1
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        set_cell(table, row_idx, 0, dim, font_size=12, bold=True,
                 color=DARK_BLUE, fill_color=bg, alignment=PP_ALIGN.CENTER)
        set_cell(table, row_idx, 1, before, font_size=12, color=RED,
                 fill_color=bg, alignment=PP_ALIGN.CENTER)
        set_cell(table, row_idx, 2, after, font_size=12, bold=True,
                 color=GREEN, fill_color=bg, alignment=PP_ALIGN.CENTER)

    # Sample report
    add_textbox(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(0.4),
                '巡检报告示例', font_size=16, bold=True, color=DARK_BLUE)

    report_box = add_rect(slide, Inches(7.5), Inches(1.9), Inches(5.3),
                          Inches(4.8), fill_color=RGBColor(0x1E, 0x1E, 0x1E))
    report_text = (
        '## 巡检报告\n'
        '事件：XX 业务访问超时\n'
        '时间：2025-01-15 14:30\n'
        '\n'
        '### 告警概况\n'
        '- [CRITICAL] HTTP 5xx 错误率 > 5%\n'
        '- [WARNING] CPU 使用率 85%\n'
        '- [WARNING] 内存使用率 90%\n'
        '\n'
        '### 指标异常\n'
        '- HTTP P99 延迟: 5.2s (正常 < 500ms)\n'
        '- 请求成功率: 94.8% (正常 > 99.5%)\n'
        '\n'
        '### 日志摘要\n'
        '- 近 1h 共 523 条 5xx 错误日志\n'
        '- 主要错误: Connection timeout\n'
        '\n'
        '### 影响评估\n'
        '发现 3 个活跃告警，服务可用性下降'
    )
    add_textbox(slide, Inches(7.7), Inches(2.0), Inches(4.9), Inches(4.5),
                report_text, font_size=10, color=RGBColor(0xCE, 0xD4, 0xDA))


def build_slide11_roadmap(slide):
    """Slide 11: Roadmap + Summary."""
    add_title_bar(slide, '实现路线图 + 总结', 11)

    # Roadmap phases
    phases = [
        ('阶段一', '框架搭建', ['项目骨架 + 配置管理', 'Prometheus 适配器', '基础巡检流程（可 Demo）'],
         BLUE),
        ('阶段二', '核心能力', ['AI 深入调查（Tool Calling）', '报告生成', 'LLM 抽象层对接'],
         ORANGE),
        ('阶段三', '系统扩展', ['Zabbix / ELK 适配器', '钉钉/飞书通知', '更多监控系统接入'],
         GREEN),
        ('阶段四', '上线优化', ['集成测试 + 压测', '文档编写', '正式上线运行'],
         DARK_BLUE),
    ]

    card_w = Inches(2.85)
    card_h = Inches(2.6)
    start_x = Inches(0.5)
    gap = Inches(0.25)
    y = Inches(1.5)

    for i, (phase, title, items, color) in enumerate(phases):
        x = start_x + i * (card_w + gap)
        card = add_rect(slide, x, y, card_w, card_h, fill_color=WHITE,
                        line_color=color)
        card.line.width = Pt(1.5)
        # Phase header
        add_rounded_rect(slide, x + Inches(0.1), y + Inches(0.1),
                         card_w - Inches(0.2), Inches(0.5),
                         fill_color=color,
                         text=f'{phase}：{title}', font_size=13)
        # Items
        for j, item in enumerate(items):
            add_textbox(slide, x + Inches(0.2), y + Inches(0.8 + j * 0.35),
                        card_w - Inches(0.4), Inches(0.32),
                        '\u2022  ' + item, font_size=11, color=BLACK)

        # Arrow between phases
        if i < len(phases) - 1:
            add_arrow_right(slide, x + card_w + Inches(0.02), y + Inches(1.0),
                            Inches(0.2), Inches(0.35), fill_color=GRAY)

    # Summary section
    summary_y = Inches(4.5)
    add_rect(slide, Inches(0.5), summary_y, Inches(12.3), Inches(2.2),
             fill_color=DARK_BLUE)

    add_textbox(slide, Inches(0.8), summary_y + Inches(0.15), Inches(11), Inches(0.5),
                '核心价值', font_size=20, bold=True, color=WHITE)

    values = [
        ('效率提升', '巡检时间从 15-30 分钟缩短到 1-2 分钟', ORANGE),
        ('全面覆盖', '自动扫描所有监控系统，不遗漏任何告警', BLUE),
        ('标准化输出', '统一格式的巡检报告，经验可沉淀', GREEN),
    ]
    for i, (title, desc, color) in enumerate(values):
        vx = Inches(0.8 + i * 4.1)
        add_rounded_rect(slide, vx, summary_y + Inches(0.7), Inches(1.5),
                         Inches(0.38), fill_color=color, text=title,
                         font_size=12, bold=True)
        add_textbox(slide, vx, summary_y + Inches(1.2), Inches(3.8), Inches(0.5),
                    desc, font_size=12, color=WHITE)

    # CTA
    add_rounded_rect(slide, Inches(4.5), summary_y + Inches(1.7), Inches(4.3),
                     Inches(0.45), fill_color=ORANGE,
                     text='请求立项审批，启动阶段一开发', font_size=15, bold=True)


# ── Main ────────────────────────────────────────────────────────────────

def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_dir = os.path.dirname(script_dir)
    template_path = os.path.join(project_dir, 'docs', 'ppt template.pptx')
    output_path = os.path.join(project_dir, 'docs', 'operation-agent-proposal.pptx')

    prs = Presentation(template_path)
    layout = prs.slide_layouts[0]

    # Delete existing empty slide
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[0]
    prs.slides._sldIdLst.remove(sldId)

    # Build all 11 slides
    builders = [
        build_slide1_cover,
        build_slide2_pain_points,
        build_slide3_solution_overview,
        build_slide4_architecture,
        build_slide5_three_phases,
        build_slide6_agent_detail,
        build_slide7_interaction,
        build_slide8_adapters,
        build_slide9_deployment,
        build_slide10_comparison,
        build_slide11_roadmap,
    ]

    for i, builder in enumerate(builders):
        slide = prs.slides.add_slide(layout)
        print(f'Building slide {i + 1}: {builder.__name__}')
        builder(slide)

    prs.save(output_path)
    print(f'\nSaved to: {output_path}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
